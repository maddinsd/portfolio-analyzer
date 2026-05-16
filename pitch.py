from __future__ import annotations

import math
from datetime import date as _date
from pathlib import Path

from pptx import Presentation
from utils import get_conviction
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── Design constants ──────────────────────────────────────────────────────────
_W, _H       = Inches(13.3), Inches(7.5)
_NAVY        = RGBColor(0x00, 0x33, 0x66)
_BLUE        = RGBColor(0x1F, 0x4E, 0x79)
_STEEL       = RGBColor(0x2D, 0x5F, 0x8A)
_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
_DGREY       = RGBColor(0x33, 0x33, 0x33)
_LGREY       = RGBColor(0xF2, 0xF2, 0xF2)
_MGREY       = RGBColor(0x70, 0x70, 0x70)
_GREEN       = RGBColor(0x1D, 0x6F, 0x42)
_RED         = RGBColor(0xC0, 0x00, 0x00)
_GOLD        = RGBColor(0xD4, 0xA0, 0x17)
_FONT        = "Calibri"
_GS_LBL      = "University of Cincinnati  |  Lindner College of Business"
_HDR_H       = Inches(0.62)
_CT          = Inches(0.72)   # content top (below header)
_ML          = Inches(0.4)    # left margin


# ── Formatting helpers ────────────────────────────────────────────────────────

def _sf(val, digits: int = 2):
    if val is None:
        return None
    try:
        f = float(val)
        return None if (math.isnan(f) or math.isinf(f)) else round(f, digits)
    except (TypeError, ValueError):
        return None


def _sa(d: dict | None, *keys, default="—"):
    """Safe nested dict access."""
    if not d:
        return default
    for k in keys:
        if isinstance(d, dict):
            d = d.get(k)
        else:
            return default
    v = d
    if v is None:
        return default
    return v


def _fmt_price(v) -> str:
    f = _sf(v)
    return f"${f:,.2f}" if f is not None else "—"


def _fmt_pct(v, plus: bool = False) -> str:
    f = _sf(v, 1)
    if f is None:
        return "—"
    prefix = "+" if plus and f > 0 else ""
    return f"{prefix}{f:.1f}%"


def _fmt_large(v) -> str:
    f = _sf(v)
    if f is None:
        return "—"
    for suffix, thr in (("T", 1e12), ("B", 1e9), ("M", 1e6)):
        if abs(f) >= thr:
            return f"${f/thr:.2f}{suffix}"
    return f"${f:,.0f}"


def _fmt_m(v_m) -> str:
    f = _sf(v_m)
    if f is None:
        return "—"
    if abs(f) >= 1000:
        return f"${f/1000:.1f}B"
    return f"${f:.0f}M"


def _rec_color(rating: str | None) -> RGBColor:
    r = (rating or "").upper()
    if r == "BUY":
        return _GREEN
    if r == "SELL":
        return _RED
    return _MGREY


# ── Data-driven fallbacks (used when research pipeline is unavailable) ────────

def _fallback_rating(research: dict | None, cov_result: dict | None) -> str:
    thesis = (research or {}).get("thesis") or {}
    if not thesis.get("_placeholder") and thesis.get("rating"):
        return thesis["rating"]
    if cov_result and not cov_result.get("error"):
        return cov_result.get("consensus_rating") or "—"
    return "—"


def _fallback_target(research: dict | None, cov_result: dict | None) -> str:
    thesis = (research or {}).get("thesis") or {}
    if not thesis.get("_placeholder"):
        raw = (thesis.get("target") or "").split("—")[0].strip()
        if raw:
            return raw
    if cov_result and not cov_result.get("error"):
        mean_t = cov_result.get("mean_target")
        if mean_t:
            return f"{_fmt_price(mean_t)} (analyst consensus)"
    return "—"


def _build_data_catalysts(stats: dict, fin_data: dict,
                          comp_result: dict | None, cov_result: dict | None) -> list[str]:
    info  = stats.get("info", {})
    a_inc = (fin_data.get("income_statement") or {}).get("annual") or {}
    cats  = []

    rev_gr      = stats.get("revenue_growth_yoy")
    peer_rev_gr = _sf(_sa(comp_result, "peer_medians", "rev_growth")) if comp_result and not comp_result.get("error") else None
    if rev_gr is not None:
        vs_peer = f" vs. peer median {peer_rev_gr:.1f}%" if peer_rev_gr else ""
        cats.append(f"Revenue {_fmt_pct(rev_gr, plus=True)} YoY{vs_peer} — sustained top-line growth at scale")

    gm_list   = a_inc.get("gross_margin") or []
    gm        = gm_list[0] if gm_list else None
    gm_prior  = gm_list[1] if len(gm_list) > 1 else None
    if gm is not None and gm_prior is not None:
        delta = round(gm - gm_prior, 1)
        cats.append(f"Gross margin {gm:.1f}% ({'+' if delta >= 0 else ''}{delta:.1f}pp YoY) — operating leverage from services mix shift driving margin re-rating")

    n_an   = cov_result.get("total_analysts") if cov_result and not cov_result.get("error") else None
    mean_t = cov_result.get("mean_target")    if cov_result and not cov_result.get("error") else None
    up     = cov_result.get("upside_pct")     if cov_result and not cov_result.get("error") else None
    rat    = cov_result.get("consensus_rating") if cov_result and not cov_result.get("error") else None
    if mean_t and n_an:
        cats.append(f"{n_an} analysts covering; {rat or 'consensus'} with mean target {_fmt_price(mean_t)} ({_fmt_pct(up, plus=True)} upside)")

    return cats[:3] or ["Run without --dry-run for analyst-specific catalyst data"]


def _build_data_risks(stats: dict, fin_data: dict,
                      dcf_result: dict | None, comp_result: dict | None) -> list[str]:
    info  = stats.get("info", {})
    risks = []

    fpe      = _sf(info.get("forwardPE"), 1)
    peer_fpe = _sf(_sa(comp_result, "peer_medians", "fpe")) if comp_result and not comp_result.get("error") else None
    if fpe and peer_fpe and peer_fpe > 0:
        prem = round((fpe / peer_fpe - 1) * 100, 0)
        risks.append(f"Valuation compression — {fpe:.1f}x fwd P/E ({prem:.0f}% premium to {peer_fpe:.1f}x peer median); any guidance miss could trigger immediate de-rating")

    iv   = _sf(dcf_result.get("valuation", {}).get("intrinsic")) if dcf_result and not dcf_result.get("error") else None
    px   = _sf(stats.get("current_price"))
    wacc = dcf_result.get("inputs", {}).get("wacc") if dcf_result and not dcf_result.get("error") else None
    if iv and px and iv < px:
        disc = round((1 - iv / px) * 100, 0)
        risks.append(f"Intrinsic value gap — DCF fair value {_fmt_price(iv)} is {disc:.0f}% below current price at WACC {wacc}%; priced for perfect 5-year execution")

    rev_gr   = stats.get("revenue_growth_yoy")
    rev_cagr = dcf_result.get("inputs", {}).get("rev_cagr") if dcf_result and not dcf_result.get("error") else None
    if rev_gr is not None:
        risks.append(f"Growth deceleration — LTM revenue {_fmt_pct(rev_gr, plus=True)} vs. DCF CAGR {_fmt_pct(rev_cagr)}; any slowdown destroys terminal value at current multiples")

    return risks[:3] or ["Run without --dry-run for detailed risk analysis"]


def _build_data_bulls(stats: dict, fin_data: dict, comp_result: dict | None) -> list[str]:
    info  = stats.get("info", {})
    a_inc = (fin_data.get("income_statement") or {}).get("annual") or {}
    bulls = []

    rev_gr = stats.get("revenue_growth_yoy")
    if rev_gr is not None:
        bulls.append(f"Revenue {_fmt_pct(rev_gr, plus=True)} YoY — top-line expansion sustained at scale demonstrates pricing power that commoditizing hardware peers cannot replicate")

    gm_list  = a_inc.get("gross_margin") or []
    gm       = gm_list[0] if gm_list else None
    gm_prior = gm_list[1] if len(gm_list) > 1 else None
    if gm is not None and gm_prior is not None:
        delta = round(gm - gm_prior, 1)
        bulls.append(f"Gross margin {gm:.1f}% ({'+' if delta >= 0 else ''}{delta:.1f}pp YoY) — services mix shift is accelerating margin expansion beyond what the market is pricing in on a hardware-only multiple")

    fpe      = _sf(info.get("forwardPE"), 1)
    peer_fpe = _sf(_sa(comp_result, "peer_medians", "fpe")) if comp_result and not comp_result.get("error") else None
    if fpe and peer_fpe and peer_fpe > 0:
        prem = round((fpe / peer_fpe - 1) * 100, 0)
        bulls.append(f"{fpe:.1f}x fwd P/E ({prem:.0f}% premium to {peer_fpe:.1f}x peer median) justified by ecosystem lock-in — switching cost is platform, not device; market underestimates duration")

    beta = _sf(info.get("beta"), 2)
    if beta:
        bulls.append(f"Beta {beta:.2f} — defensive large-cap quality with deep institutional ownership provides downside support; outperforms in risk-off vs. higher-beta sector peers")

    return bulls[:4] or ["Run without --dry-run for AI-generated investment thesis"]


# ── Drawing helpers ───────────────────────────────────────────────────────────

def _rect(slide, left, top, width, height,
          fill: RGBColor = _NAVY, border: RGBColor | None = None):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _SHP
    shp = slide.shapes.add_shape(_SHP.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def _txt(slide, text: str, left, top, width, height, *,
         size: int = 12, bold: bool = False, italic: bool = False,
         color: RGBColor = _DGREY, align=PP_ALIGN.LEFT,
         valign=MSO_ANCHOR.TOP, wrap: bool = True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = _FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _txt2(slide, lines: list[tuple[str, dict]], left, top, width, height, *,
          wrap: bool = True):
    """Multi-paragraph textbox. lines = [(text, {size, bold, italic, color, align}), ...]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        if opts.get("space_before"):
            p.space_before = Pt(opts["space_before"])
        run = p.add_run()
        run.text = text
        run.font.name = _FONT
        run.font.size = Pt(opts.get("size", 12))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", _DGREY)
    return tb


def _add_header(slide, title: str, page_num: int):
    _rect(slide, 0, 0, _W, _HDR_H, fill=_NAVY)
    _txt(slide, title.upper(), _ML, Inches(0.12), Inches(9.5), Inches(0.42),
         size=16, bold=True, color=_WHITE, align=PP_ALIGN.LEFT,
         valign=MSO_ANCHOR.MIDDLE)
    _txt(slide, _GS_LBL, Inches(9.5), Inches(0.12), Inches(3.5), Inches(0.42),
         size=10, italic=True, color=_WHITE, align=PP_ALIGN.RIGHT,
         valign=MSO_ANCHOR.MIDDLE)
    _txt(slide, str(page_num), Inches(12.7), Inches(7.2), Inches(0.35), Inches(0.2),
         size=9, color=_MGREY, align=PP_ALIGN.RIGHT)


def _cell(table, row: int, col: int, text: str, *,
          size: int = 11, bold: bool = False, color: RGBColor = _DGREY,
          fill: RGBColor | None = None, align=PP_ALIGN.LEFT):
    cell = table.cell(row, col)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text)
    run.font.name = _FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


# ── Football field data extraction ────────────────────────────────────────────

def _football_data(stats: dict, dcf_result: dict | None,
                   comp_result: dict | None, cov_result: dict | None) -> dict:
    info = stats.get("info", {})
    px   = _sf(stats.get("current_price"))

    # DCF range: sensitivity table corners [WACC+2%,g=1.5%] to [WACC-2%,g=3.5%]
    dcf_low = dcf_high = None
    if dcf_result and not dcf_result.get("error"):
        tbl = dcf_result.get("sensitivity", {}).get("table", [])
        if tbl and len(tbl) >= 5:
            dcf_low  = _sf(tbl[4][0])   # WACC+2%, lowest growth
            dcf_high = _sf(tbl[0][4])   # WACC-2%, highest growth

    # 52-week range
    wk52_low  = _sf(info.get("fiftyTwoWeekLow"))
    wk52_high = _sf(info.get("fiftyTwoWeekHigh"))

    # Analyst targets — individual targets preferred; fall back to mean ± spread estimate
    anal_low = anal_high = None
    if cov_result and not cov_result.get("error"):
        anal_low  = _sf(cov_result.get("low_target"))
        anal_high = _sf(cov_result.get("high_target"))
        if anal_low is None or anal_high is None:
            mean_t = _sf(cov_result.get("mean_target"))
            if mean_t:
                spread = _sf(cov_result.get("target_spread_pct")) or 20.0
                anal_low  = round(mean_t * (1 - spread / 200), 2)
                anal_high = round(mean_t * (1 + spread / 200), 2)

    # Comps-implied: peer fpe range × target forward EPS
    comp_low = comp_high = None
    eps = _sf(info.get("forwardEps")) or _sf(info.get("trailingEps"))
    if eps and eps > 0 and comp_result and not comp_result.get("error"):
        peer_fpes = sorted(_sf(p.get("fpe")) for p in comp_result.get("peers", [])
                           if _sf(p.get("fpe")) and _sf(p.get("fpe")) > 0)
        if peer_fpes:
            def _q(p, s=peer_fpes):
                k = (len(s) - 1) * p
                lo, hi = int(k), min(int(k) + 1, len(s) - 1)
                return s[lo] + (s[hi] - s[lo]) * (k - lo)
            comp_low  = round(_q(0.25) * eps, 2)
            comp_high = round(_q(0.75) * eps, 2)

    bars = [
        ("DCF",            dcf_low,   dcf_high,   _NAVY),
        ("Comps-Implied",  comp_low,  comp_high,  _STEEL),
        ("52-Week Range",  wk52_low,  wk52_high,  _BLUE),
        ("Analyst Targets", anal_low, anal_high,  _GOLD),
    ]

    all_vals = [v for b in bars for v in (b[1], b[2]) if v] + ([px] if px else [])
    if not all_vals:
        return {"bars": bars, "chart_min": 0, "chart_max": 1, "current_price": px}

    chart_min = min(all_vals) * 0.90
    chart_max = max(all_vals) * 1.10
    return {"bars": bars, "chart_min": chart_min, "chart_max": chart_max, "current_price": px}


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(prs, stats: dict, research: dict | None, cov_result: dict | None = None):
    info  = stats.get("info", {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, _W, _H, fill=_NAVY)

    company = info.get("shortName") or info.get("longName") or stats.get("ticker", "")
    ticker  = stats.get("ticker", "")
    sector  = info.get("sector") or "—"
    exch    = info.get("fmp_exchange") or info.get("exchange") or "NASDAQ"
    px      = _fmt_price(stats.get("current_price"))
    mktcap  = _fmt_large(info.get("marketCap"))

    rating = _fallback_rating(research, cov_result)
    target = _fallback_target(research, cov_result)

    _txt(slide, _GS_LBL, Inches(9.5), Inches(0.15), Inches(3.5), Inches(0.3),
         size=10, italic=True, color=_MGREY, align=PP_ALIGN.RIGHT)

    _txt(slide, company, Inches(1.0), Inches(1.1), Inches(11.3), Inches(1.0),
         size=44, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    _txt(slide, f"{ticker}  |  {exch}  |  {sector}",
         Inches(1.0), Inches(2.1), Inches(11.3), Inches(0.4),
         size=16, color=_MGREY, align=PP_ALIGN.CENTER)

    rc = _rec_color(rating)
    _rect(slide, Inches(5.15), Inches(2.7), Inches(3.0), Inches(0.72), fill=rc)
    _txt(slide, rating if rating != "—" else "UNDER REVIEW",
         Inches(5.15), Inches(2.7), Inches(3.0), Inches(0.72),
         size=32, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _txt(slide, f"Target Price: {target}",
         Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.4),
         size=20, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    _txt(slide, f"Current Price: {px}   |   Market Cap: {mktcap}",
         Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.35),
         size=14, color=_LGREY, align=PP_ALIGN.CENTER)

    _txt(slide, str(_date.today().strftime("%B %d, %Y")),
         Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.3),
         size=12, color=_MGREY, align=PP_ALIGN.CENTER)

    _txt(slide, "Analyst: Samuel Madding  |  University of Cincinnati | Lindner College of Business — For Educational Purposes Only",
         Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.25),
         size=9, italic=True, color=_MGREY, align=PP_ALIGN.CENTER)


def _slide_summary(prs, stats: dict, fin_data: dict,
                   research: dict | None, comp_result: dict | None,
                   cov_result: dict | None, dcf_result: dict | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Investment Summary", 2)

    thesis  = (research or {}).get("thesis") or {}
    ok      = not thesis.get("_placeholder") and bool(thesis.get("rating") or thesis.get("catalysts"))
    rating  = _fallback_rating(research, cov_result)
    target  = _fallback_target(research, cov_result)
    cats    = thesis.get("catalysts", []) if ok else []
    bears   = thesis.get("bear", [])     if ok else []
    verdict = thesis.get("verdict", "")  if ok else ""

    if not cats:
        cats = _build_data_catalysts(stats, fin_data, comp_result, cov_result)
    if not bears:
        bears = _build_data_risks(stats, fin_data, dcf_result, comp_result)

    px     = stats.get("current_price")
    upside = cov_result.get("upside_pct") if cov_result and not cov_result.get("error") else None

    # Left panel — rating card
    rc = _rec_color(rating)
    _rect(slide, _ML, _CT, Inches(4.0), Inches(1.0), fill=rc)
    _txt(slide, rating if rating != "—" else "COVERAGE",
         _ML, _CT, Inches(4.0), Inches(1.0),
         size=36, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    up_color = _GREEN if (upside or 0) >= 0 else _RED
    _txt2(slide, [
        (f"Target: {target}", {"size": 16, "bold": True, "color": _NAVY}),
        (f"Current: {_fmt_price(px)}", {"size": 13, "color": _DGREY, "space_before": 4}),
        (f"Upside: {_fmt_pct(upside, plus=True)}", {"size": 14, "bold": True, "color": up_color, "space_before": 4}),
        ("12-Month Horizon", {"size": 10, "color": _MGREY, "italic": True, "space_before": 6}),
    ], _ML, Inches(1.82), Inches(4.0), Inches(1.8))

    # Valuation context
    tgt_fpe  = _sf(stats.get("info", {}).get("forwardPE"), 1)
    peer_fpe = _sf(_sa(comp_result, "peer_medians", "fpe")) if comp_result and not comp_result.get("error") else None
    if tgt_fpe and peer_fpe and peer_fpe > 0:
        prem = round((tgt_fpe / peer_fpe - 1) * 100, 1)
        sign = "premium" if prem >= 0 else "discount"
        val_line = f"Trading at {tgt_fpe:.1f}x fwd P/E vs. peer median {peer_fpe:.1f}x — {abs(prem):.1f}% {sign}"
    else:
        val_line = ""
    if val_line:
        _rect(slide, _ML, Inches(3.7), Inches(4.0), Inches(0.5), fill=_LGREY)
        _txt(slide, val_line, _ML + Inches(0.1), Inches(3.75), Inches(3.85), Inches(0.45),
             size=10, italic=True, color=_NAVY, valign=MSO_ANCHOR.MIDDLE)

    # Right panel — catalysts + risks
    _txt(slide, "Key Catalysts", Inches(4.8), _CT, Inches(8.1), Inches(0.35),
         size=14, bold=True, color=_NAVY)
    cat_lines = [(f"▸  {c}", {"size": 12, "color": _DGREY, "space_before": 5}) for c in cats[:3]]
    _txt2(slide, cat_lines, Inches(4.8), Inches(1.15), Inches(8.1), Inches(2.0))

    _rect(slide, Inches(4.8), Inches(3.3), Inches(8.1), Inches(0.02), fill=_LGREY)

    _txt(slide, "Key Risks", Inches(4.8), Inches(3.4), Inches(8.1), Inches(0.3),
         size=13, bold=True, color=_RED)
    risk_lines = [(f"▸  {b}", {"size": 11, "color": _DGREY, "space_before": 4}) for b in bears[:2]]
    _txt2(slide, risk_lines, Inches(4.8), Inches(3.75), Inches(8.1), Inches(1.4))

    if verdict:
        _rect(slide, 0, Inches(5.5), _W, Inches(0.75), fill=_BLUE)
        _txt(slide, f'"{verdict}"',
             _ML, Inches(5.55), Inches(12.9), Inches(0.65),
             size=11, italic=True, color=_WHITE, valign=MSO_ANCHOR.MIDDLE)


def _slide_overview(prs, stats: dict, fin_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Company Overview", 3)

    info     = stats.get("info", {})
    desc_raw = info.get("longBusinessSummary") or "Business description not available."
    _txt(slide, desc_raw, _ML, _CT, Inches(12.5), Inches(1.35),
         size=9.5, color=_DGREY, wrap=True)

    _rect(slide, _ML, Inches(2.15), Inches(12.5), Inches(0.02), fill=_LGREY)

    # Left column: profile stats
    ceo      = info.get("fmp_ceo") or "—"
    emp      = info.get("fullTimeEmployees")
    emp_str  = f"{emp:,}" if emp else "—"
    city     = info.get("fmp_city") or info.get("city") or ""
    state    = info.get("fmp_state") or info.get("state") or ""
    hq       = f"{city}, {state}".strip(", ") or "—"
    ipo_raw  = (info.get("fmp_ipo_date") or "")[:10]
    try:
        from datetime import datetime as _dt
        ipo = _dt.strptime(ipo_raw, "%Y-%m-%d").strftime("%b %d, %Y") if ipo_raw else "—"
    except Exception:
        ipo = ipo_raw or "—"
    sector   = info.get("sector") or "—"
    industry = info.get("industry") or "—"

    left_data = [
        ("CEO",       ceo),
        ("Employees", emp_str),
        ("HQ",        hq),
        ("IPO Date",  ipo),
        ("Sector",    sector),
        ("Industry",  industry),
    ]
    y = Inches(2.25)
    for label, value in left_data:
        _txt(slide, label, _ML, y, Inches(1.6), Inches(0.32),
             size=10, bold=True, color=_NAVY)
        _txt(slide, value, Inches(2.1), y, Inches(3.8), Inches(0.32),
             size=10, color=_DGREY)
        y += Inches(0.38)

    # Right column: key market metrics
    px       = _fmt_price(stats.get("current_price"))
    mktcap   = _fmt_large(info.get("marketCap"))
    hi52     = _fmt_price(info.get("fiftyTwoWeekHigh"))
    lo52     = _fmt_price(info.get("fiftyTwoWeekLow"))
    pe       = _sf(info.get("trailingPE"), 1)
    fpe      = _sf(info.get("forwardPE"), 1)
    rev_gr   = _fmt_pct(stats.get("revenue_growth_yoy"), plus=True)
    beta     = _sf(info.get("beta"), 2)
    dy_f     = _sf(info.get("dividendYield"), 4)
    div_yld  = f"{dy_f:.2f}%" if dy_f is not None else "—"

    right_data = [
        ("Price",         px),
        ("Market Cap",    mktcap),
        ("52-Wk Range",   f"{lo52} – {hi52}"),
        ("P/E (TTM)",     f"{pe}x" if pe else "—"),
        ("P/E (Fwd)",     f"{fpe}x" if fpe else "—"),
        ("Revenue Growth", rev_gr),
        ("Beta",          str(beta) if beta else "—"),
        ("Div. Yield",    div_yld),
    ]
    y = Inches(2.25)
    for label, value in right_data:
        _txt(slide, label, Inches(6.5), y, Inches(2.2), Inches(0.32),
             size=10, bold=True, color=_NAVY)
        _txt(slide, value, Inches(8.8), y, Inches(4.0), Inches(0.32),
             size=10, color=_DGREY)
        y += Inches(0.38)

    # ── Financial highlights strip (fills bottom of slide) ────────────────────
    a_inc = (fin_data.get("income_statement") or {}).get("annual") or {}
    a_cf  = (fin_data.get("cash_flow") or {}).get("annual") or {}
    rev_list = a_inc.get("revenue", [])
    gm_list  = a_inc.get("gross_margin", [])
    om_list  = a_inc.get("operating_margin", [])
    fcf_list = a_cf.get("free_cash_flow", [])

    def _fmt_rev(v_m):
        if v_m is None: return "—"
        f = _sf(v_m)
        return f"${f/1000:.1f}B" if abs(f) >= 1000 else f"${f:.0f}M"

    def _fmt_margin(v):
        f = _sf(v, 1)
        return f"{f:.1f}%" if f is not None else "—"

    fin_highlights = [
        ("LTM Revenue",    _fmt_rev(rev_list[0] if rev_list else None)),
        ("Gross Margin",   _fmt_margin(gm_list[0] if gm_list else None)),
        ("Op. Margin",     _fmt_margin(om_list[0] if om_list else None)),
        ("Free Cash Flow", _fmt_rev(fcf_list[0] if fcf_list else None)),
    ]

    strip_y = Inches(5.1)
    strip_h = Inches(1.8)
    col_w   = Inches(12.5) / len(fin_highlights)
    _rect(slide, _ML, strip_y, Inches(12.5), Inches(0.02), fill=_LGREY)
    _txt(slide, "FINANCIAL HIGHLIGHTS",
         _ML, strip_y + Inches(0.06), Inches(12.5), Inches(0.25),
         size=8, bold=True, color=_MGREY)
    for i, (lbl, val) in enumerate(fin_highlights):
        cx = _ML + col_w * i
        _txt(slide, val,   cx, strip_y + Inches(0.36), col_w, Inches(0.5),
             size=20, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
        _txt(slide, lbl,   cx, strip_y + Inches(0.9),  col_w, Inches(0.3),
             size=9, color=_MGREY, align=PP_ALIGN.CENTER)


def _slide_thesis(prs, research: dict | None,
                  stats: dict | None = None, fin_data: dict | None = None,
                  comp_result: dict | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Investment Thesis", 4)

    thesis  = (research or {}).get("thesis") or {}
    ok      = not thesis.get("_placeholder")
    bulls   = thesis.get("bull", []) if ok else []
    verdict = thesis.get("verdict", "") if ok else ""

    if len(bulls) < 2 and stats and fin_data:
        points = _build_data_bulls(stats, fin_data, comp_result)
    else:
        points = bulls[:4]

    BOX_H = Inches(2.15)
    ROW2_Y = Inches(3.02)   # tighter gap between rows
    positions = [
        (_ML,            _CT),
        (Inches(6.85),   _CT),
        (_ML,            ROW2_Y),
        (Inches(6.85),   ROW2_Y),
    ]
    for i, (text, (cx, cy)) in enumerate(zip(points[:4], positions)):
        _rect(slide, cx, cy, Inches(5.9), BOX_H, fill=_LGREY, border=_LGREY)
        _rect(slide, cx, cy, Inches(0.06), BOX_H, fill=_NAVY)
        _txt(slide, text,
             cx + Inches(0.15), cy + Inches(0.1),
             Inches(5.65), BOX_H - Inches(0.12),
             size=8.5, color=_DGREY, wrap=True)

    if verdict:
        _rect(slide, 0, Inches(5.38), _W, Inches(0.6), fill=_NAVY)
        _txt(slide, f'Analyst View: "{verdict}"',
             _ML, Inches(5.43), Inches(12.9), Inches(0.5),
             size=11, italic=True, color=_WHITE, valign=MSO_ANCHOR.MIDDLE)


def _slide_financials(prs, stats: dict, fin_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Financial Snapshot", 5)

    a_inc = (fin_data.get("income_statement") or {}).get("annual") or {}
    a_cf  = (fin_data.get("cash_flow") or {}).get("annual") or {}
    bs    = fin_data.get("balance_sheet") or {}

    dates   = a_inc.get("dates", [])[:3]
    rev     = a_inc.get("revenue", [None]*3)[:3]
    gm      = a_inc.get("gross_margin", [None]*3)[:3]
    om      = a_inc.get("operating_margin", [None]*3)[:3]
    nm      = a_inc.get("net_margin", [None]*3)[:3]
    fcf     = a_cf.get("free_cash_flow", [None]*3)[:3]
    yoy_rev = a_inc.get("yoy_revenue", [None]*3)[:3]

    n_cols  = 1 + len(dates) + 1   # Metric + N date cols + YoY
    n_rows  = 6                    # header + 5 data rows
    tbl_shp = slide.shapes.add_table(
        n_rows, n_cols, _ML, _CT, Inches(12.5), Inches(2.85))
    tbl     = tbl_shp.table

    # Column widths
    tbl.columns[0].width = Inches(2.5)
    date_w = Inches(2.5) if len(dates) >= 3 else Inches(3.0)
    for c in range(1, n_cols - 1):
        tbl.columns[c].width = date_w
    tbl.columns[n_cols - 1].width = Inches(1.8)

    # Header row
    _cell(tbl, 0, 0, "Metric", bold=True, color=_WHITE, fill=_NAVY, size=11, align=PP_ALIGN.LEFT)
    for c, d in enumerate(dates, 1):
        _cell(tbl, 0, c, d, bold=True, color=_WHITE, fill=_NAVY, size=11, align=PP_ALIGN.CENTER)
    _cell(tbl, 0, n_cols - 1, "YoY", bold=True, color=_WHITE, fill=_NAVY, size=11, align=PP_ALIGN.CENTER)

    def _pct_str(v):
        return f"{v:.1f}%" if v is not None else "—"

    rows_data = [
        ("Revenue",          [_fmt_m(v) for v in rev],  _fmt_pct(yoy_rev[0] if yoy_rev else None, plus=True)),
        ("Gross Margin",     [_pct_str(v) for v in gm], "—"),
        ("Operating Margin", [_pct_str(v) for v in om], "—"),
        ("Net Margin",       [_pct_str(v) for v in nm], "—"),
        ("Free Cash Flow",   [_fmt_m(v) for v in fcf],  _fmt_pct(a_cf.get("yoy_fcf"), plus=True)),
    ]

    for r, (label, vals, yoy_str) in enumerate(rows_data, 1):
        row_fill = _LGREY if r % 2 == 0 else _WHITE
        _cell(tbl, r, 0, label, bold=True, color=_NAVY, fill=row_fill, size=11)
        for c, v in enumerate(vals, 1):
            _cell(tbl, r, c, v, color=_DGREY, fill=row_fill, size=11, align=PP_ALIGN.CENTER)
        yoy_color = _GREEN if "+" in yoy_str else (_RED if "-" in yoy_str and yoy_str != "—" else _MGREY)
        _cell(tbl, r, n_cols - 1, yoy_str, bold=True, color=yoy_color, fill=row_fill,
              size=11, align=PP_ALIGN.CENTER)

    # Bottom callout boxes
    info      = stats.get("info", {})
    callouts  = [
        ("Market Cap",    _fmt_large(info.get("marketCap"))),
        ("Net Debt",      _fmt_m(bs.get("net_debt"))),
        ("P/E (TTM)",     f"{_sf(info.get('trailingPE'),1):.1f}x" if _sf(info.get("trailingPE")) else "—"),
        ("P/E (Fwd)",     f"{_sf(info.get('forwardPE'),1):.1f}x"  if _sf(info.get("forwardPE"))  else "—"),
    ]
    box_w = Inches(2.9)
    for i, (label, value) in enumerate(callouts):
        cx = _ML + Inches(i * 3.1)
        _rect(slide, cx, Inches(4.0), box_w, Inches(0.95), fill=_NAVY)
        _txt(slide, value, cx, Inches(4.05), box_w, Inches(0.5),
             size=22, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _txt(slide, label, cx, Inches(4.6), box_w, Inches(0.3),
             size=10, color=_LGREY, align=PP_ALIGN.CENTER)


def _slide_dcf(prs, stats: dict, fin_data: dict, dcf_result: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "DCF Valuation", 6)

    if not dcf_result or dcf_result.get("error"):
        err = (dcf_result or {}).get("error", "DCF data unavailable")
        _txt(slide, f"DCF analysis unavailable: {err}",
             _ML, _CT, Inches(12.5), Inches(1.0), size=12, color=_MGREY, italic=True)
        return

    inp = dcf_result.get("inputs", {})
    val = dcf_result.get("valuation", {})
    sens = dcf_result.get("sensitivity", {})

    wacc    = inp.get("wacc")
    tg      = inp.get("tg")
    beta    = inp.get("beta")
    iv      = val.get("intrinsic")
    px      = val.get("current_price") or stats.get("current_price")
    upside  = val.get("upside_pct")
    ev_m    = val.get("ev_m")
    tv_pct  = val.get("tv_pct")

    # Left panel: inputs
    _txt(slide, "DCF Inputs", _ML, _CT, Inches(3.8), Inches(0.3),
         size=13, bold=True, color=_NAVY)
    left_items = [
        ("WACC",           f"{wacc:.2f}%" if wacc else "—"),
        ("Terminal Growth", f"{tg:.1f}%"  if tg   else "—"),
        ("Beta",           f"{beta:.2f}"  if beta  else "—"),
        ("Revenue CAGR",   f"{inp.get('rev_cagr', '—'):.1f}%" if inp.get("rev_cagr") else "—"),
        ("EBIT Margin",    f"{inp.get('ebit_margin', '—'):.1f}%" if inp.get("ebit_margin") else "—"),
        ("Forecast Years", "5"),
    ]
    y = Inches(1.1)
    for label, value in left_items:
        _txt(slide, label, _ML, y, Inches(2.0), Inches(0.32), size=10, bold=True, color=_NAVY)
        _txt(slide, value, Inches(2.5), y, Inches(1.4), Inches(0.32), size=10, color=_DGREY)
        y += Inches(0.37)

    # Center panel: valuation result
    _rect(slide, Inches(4.2), _CT, Inches(4.6), Inches(2.2), fill=_LGREY, border=_LGREY)
    _txt(slide, "Intrinsic Value Per Share",
         Inches(4.2), _CT + Inches(0.1), Inches(4.6), Inches(0.35),
         size=11, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
    _txt(slide, _fmt_price(iv),
         Inches(4.2), Inches(1.2), Inches(4.6), Inches(0.9),
         size=48, bold=True, color=_NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    upside_color = _GREEN if (upside or 0) >= 0 else _RED
    _txt(slide, f"Current: {_fmt_price(px)}   |   Upside: {_fmt_pct(upside, plus=True)}",
         Inches(4.2), Inches(2.2), Inches(4.6), Inches(0.35),
         size=12, bold=True, color=upside_color, align=PP_ALIGN.CENTER)

    ev_str  = _fmt_m(ev_m)
    tv_str  = f"{tv_pct:.1f}%" if tv_pct else "—"
    _txt(slide, f"Enterprise Value: {ev_str}   |   Terminal Value: {tv_str} of EV",
         Inches(4.2), Inches(2.65), Inches(4.6), Inches(0.3),
         size=10, color=_MGREY, align=PP_ALIGN.CENTER)

    # Right panel: sensitivity table (3x3 center extract from 5x5)
    tbl_data = sens.get("table", [])
    wacc_rng = sens.get("wacc_range", [])
    tg_rng   = sens.get("tg_range", [])

    _txt(slide, "Sensitivity: WACC × Terminal Growth",
         Inches(9.2), _CT, Inches(3.8), Inches(0.3),
         size=11, bold=True, color=_NAVY)

    if tbl_data and len(tbl_data) >= 5:
        rows_idx = [0, 2, 4]
        cols_idx = [0, 2, 4]
        n_s = 4
        s_shp = slide.shapes.add_table(n_s, n_s, Inches(9.1), Inches(1.1),
                                       Inches(3.9), Inches(1.6))
        s_tbl = s_shp.table
        for c in range(n_s):
            s_tbl.columns[c].width = Inches(3.9 / n_s)

        # Header row: tg values
        _cell(s_tbl, 0, 0, "WACC \\ g", bold=True, color=_WHITE, fill=_NAVY, size=9,
              align=PP_ALIGN.CENTER)
        for ci, gi in enumerate(cols_idx, 1):
            tg_val = tg_rng[gi] if gi < len(tg_rng) else "—"
            _cell(s_tbl, 0, ci, f"{tg_val}%", bold=True, color=_WHITE, fill=_NAVY,
                  size=9, align=PP_ALIGN.CENTER)

        for ri_pos, ri in enumerate(rows_idx, 1):
            wacc_val = wacc_rng[ri] if ri < len(wacc_rng) else "—"
            row_fill = _LGREY if ri_pos % 2 == 0 else _WHITE
            _cell(s_tbl, ri_pos, 0, f"{wacc_val}%", bold=True, color=_NAVY,
                  fill=row_fill, size=9, align=PP_ALIGN.CENTER)
            for ci, gi in enumerate(cols_idx, 1):
                cell_val = tbl_data[ri][gi] if ri < len(tbl_data) and gi < len(tbl_data[ri]) else None
                is_center = (ri_pos == 2 and ci == 2)
                c_fill = _STEEL if is_center else row_fill
                c_color = _WHITE if is_center else _DGREY
                _cell(s_tbl, ri_pos, ci, _fmt_price(cell_val),
                      fill=c_fill, color=c_color, size=9, align=PP_ALIGN.CENTER)


def _slide_comps(prs, stats: dict, research: dict | None, comp_result: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Comparable Companies", 7)

    comps_r   = (research or {}).get("comps") or {}
    ok_r      = not comps_r.get("_placeholder")
    comp_list = comps_r.get("comps", []) if ok_r else []
    summary   = comps_r.get("summary", "") if ok_r else ""
    premium   = comps_r.get("premium", "") if ok_r else ""

    # Fallback: populate from comp_result["peers"] when research comps unavailable
    if not comp_list and comp_result and not comp_result.get("error"):
        for peer in comp_result.get("peers", [])[:5]:
            comp_list.append({
                "company":  peer.get("name", peer.get("ticker", "—")),
                "ticker":   peer.get("ticker", "—"),
                "ev_ebitda": None,
                "pe_fwd":   peer.get("fpe"),
                "ev_rev":   None,
            })

    ticker  = stats.get("ticker", "")
    info    = stats.get("info", {})
    tgt_pe  = _sf(info.get("trailingPE"), 1)
    tgt_fpe = _sf(info.get("forwardPE"), 1)

    n_rows  = min(len(comp_list), 5) + 2  # header + target + peers
    n_cols  = 5
    if n_rows < 2:
        n_rows = 3

    tbl_shp = slide.shapes.add_table(n_rows, n_cols, _ML, _CT, Inches(12.5), Inches(3.2))
    tbl     = tbl_shp.table
    col_ws  = [Inches(3.5), Inches(1.2), Inches(2.5), Inches(2.5), Inches(2.3)]
    for c, w in enumerate(col_ws):
        tbl.columns[c].width = w

    headers = ["Company", "Ticker", "EV/EBITDA", "P/E (Fwd)", "EV/Revenue"]
    for c, h in enumerate(headers):
        _cell(tbl, 0, c, h, bold=True, color=_WHITE, fill=_NAVY, size=11,
              align=PP_ALIGN.CENTER)

    def _val(v, suffix="x"):
        f = _sf(v, 1)
        return f"{f:.1f}{suffix}" if f else "—"

    # Target row
    _cell(tbl, 1, 0, f"{info.get('shortName', ticker)} ◀", bold=True, color=_WHITE, fill=_STEEL, size=11)
    _cell(tbl, 1, 1, ticker, bold=True, color=_WHITE, fill=_STEEL, size=11, align=PP_ALIGN.CENTER)
    _cell(tbl, 1, 2, "—", bold=True, color=_WHITE, fill=_STEEL, size=11, align=PP_ALIGN.CENTER)
    _cell(tbl, 1, 3, _val(tgt_fpe), bold=True, color=_WHITE, fill=_STEEL, size=11, align=PP_ALIGN.CENTER)
    _cell(tbl, 1, 4, "—", bold=True, color=_WHITE, fill=_STEEL, size=11, align=PP_ALIGN.CENTER)

    for r, comp in enumerate(comp_list[:n_rows - 2], 2):
        row_fill = _LGREY if r % 2 == 0 else _WHITE
        _cell(tbl, r, 0, comp.get("company", "—"), color=_DGREY, fill=row_fill, size=11)
        _cell(tbl, r, 1, comp.get("ticker", "—"), color=_DGREY, fill=row_fill, size=11, align=PP_ALIGN.CENTER)
        _cell(tbl, r, 2, _val(comp.get("ev_ebitda")), color=_DGREY, fill=row_fill, size=11, align=PP_ALIGN.CENTER)
        _cell(tbl, r, 3, _val(comp.get("pe_fwd")), color=_DGREY, fill=row_fill, size=11, align=PP_ALIGN.CENTER)
        _cell(tbl, r, 4, _val(comp.get("ev_rev")), color=_DGREY, fill=row_fill, size=11, align=PP_ALIGN.CENTER)

    # Summary
    y_sum = _CT + Inches(3.35)
    if summary:
        _txt(slide, summary, _ML, y_sum, Inches(12.5), Inches(0.7),
             size=11, color=_DGREY, wrap=True)
        y_sum += Inches(0.75)
    if premium:
        _rect(slide, _ML, y_sum, Inches(12.5), Inches(0.42), fill=_LGREY)
        _txt(slide, f"Relative Valuation: {premium}",
             _ML + Inches(0.1), y_sum + Inches(0.05), Inches(12.3), Inches(0.35),
             size=11, italic=True, color=_NAVY, valign=MSO_ANCHOR.MIDDLE)


def _slide_football(prs, stats: dict, dcf_result: dict | None,
                    comp_result: dict | None, cov_result: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Valuation Range", 8)

    fd = _football_data(stats, dcf_result, comp_result, cov_result)
    bars      = fd["bars"]
    chart_min = fd["chart_min"]
    chart_max = fd["chart_max"]
    px        = fd["current_price"]

    span  = chart_max - chart_min
    if span <= 0:
        _txt(slide, "Insufficient valuation data for football field chart.",
             _ML, _CT, Inches(12.5), Inches(1.0), size=12, color=_MGREY, italic=True)
        return

    # Chart area geometry
    label_w  = Inches(1.8)
    chart_x  = _ML + label_w
    chart_y  = _CT
    chart_w  = Inches(10.8)
    chart_h  = Inches(4.3)
    bar_h    = Inches(0.52)
    n_bars   = len(bars)
    slot_h   = chart_h / n_bars

    # Background
    _rect(slide, chart_x, chart_y, chart_w, chart_h, fill=_LGREY, border=_LGREY)

    def _x_pos(val):
        return chart_x + (val - chart_min) / span * chart_w

    # Draw bars
    colors = [_NAVY, _STEEL, _BLUE, _GOLD]
    for i, (label, low, high, color) in enumerate(bars):
        bar_cy = chart_y + i * slot_h
        bar_mid_y = bar_cy + (slot_h - bar_h) / 2

        # Methodology label
        _txt(slide, label, _ML, bar_mid_y, label_w - Inches(0.1), bar_h,
             size=10, bold=True, color=_NAVY, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)

        if low is not None and high is not None and low < high:
            x_lo = _x_pos(low)
            x_hi = _x_pos(high)
            bar_actual_w = x_hi - x_lo
            if bar_actual_w > 0:
                _rect(slide, x_lo, bar_mid_y, bar_actual_w, bar_h, fill=color)
                # Value labels — offset enough to avoid overlap with bar labels
                _txt(slide, _fmt_price(low),
                     x_lo - Inches(0.85), bar_mid_y, Inches(0.8), bar_h,
                     size=8, color=_DGREY, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
                _txt(slide, _fmt_price(high),
                     x_hi + Inches(0.04), bar_mid_y, Inches(0.8), bar_h,
                     size=8, color=_DGREY, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
        else:
            _txt(slide, "N/A", chart_x + Inches(0.1), bar_mid_y, Inches(2.0), bar_h,
                 size=10, color=_MGREY, italic=True, valign=MSO_ANCHOR.MIDDLE)

    # Current price vertical line
    if px and chart_min < px < chart_max:
        px_x = _x_pos(px)
        _rect(slide, px_x - Inches(0.025), chart_y, Inches(0.05), chart_h, fill=_RED)
        _txt(slide, f"▲ {_fmt_price(px)}",
             px_x - Inches(0.6), chart_y + chart_h + Inches(0.05), Inches(1.2), Inches(0.3),
             size=10, bold=True, color=_RED, align=PP_ALIGN.CENTER)

    # X-axis price ticks (5 marks)
    n_ticks = 5
    for t in range(n_ticks):
        tick_val = chart_min + span * t / (n_ticks - 1)
        tick_x   = _x_pos(tick_val)
        _txt(slide, _fmt_price(tick_val),
             tick_x - Inches(0.4), chart_y + chart_h + Inches(0.05),
             Inches(0.8), Inches(0.25),
             size=8, color=_MGREY, align=PP_ALIGN.CENTER)


def _slide_competitive(prs, stats: dict, comp_result: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Competitive Position", 9)

    ok = comp_result and not comp_result.get("error")
    moat = comp_result.get("claude") if ok else None
    peers = comp_result.get("peers", [])[:3] if ok else []
    medians = comp_result.get("peer_medians", {}) if ok else {}
    ranks = comp_result.get("rankings", {}) if ok else {}
    tgt = comp_result.get("target", {}) if ok else {}

    # Left: moat / ranking summary
    _txt(slide, "Competitive Assessment", _ML, _CT, Inches(6.0), Inches(0.3),
         size=13, bold=True, color=_NAVY)

    if moat and isinstance(moat, str):
        moat_text = moat[:500]
    elif ok:
        prem = ranks.get("fpe_premium_pct")
        prem_str = (f"Trading at {abs(prem):.1f}% {'premium' if prem and prem > 0 else 'discount'} "
                    f"to peer median fwd P/E. " if prem is not None else "")
        rev_rank = ranks.get("rev_growth", "mid")
        gm_rank  = ranks.get("gross_margin", "mid")
        rank_map = {"top": "above-median", "mid": "in-line with", "bot": "below-median"}
        moat_text = (f"{prem_str}Revenue growth is {rank_map[rev_rank]} peers. "
                     f"Gross margin is {rank_map[gm_rank]} peers.")
    else:
        moat_text = "Competitive data unavailable."

    _txt(slide, moat_text, _ML, Inches(1.1), Inches(6.0), Inches(1.6),
         size=11, color=_DGREY, wrap=True)

    # Peer comparison table
    _txt(slide, "Peer Comparison", _ML, Inches(2.85), Inches(6.0), Inches(0.3),
         size=12, bold=True, color=_NAVY)

    if ok and (peers or tgt):
        metrics = ["rev_growth", "gross_margin", "op_margin"]
        labels  = ["Rev Growth %", "Gross Margin %", "Op Margin %"]
        rows    = 2 + len(peers)
        cols    = 1 + len(metrics)
        p_shp   = slide.shapes.add_table(rows, cols, _ML, Inches(3.2), Inches(6.0), Inches(2.0))
        p_tbl   = p_shp.table
        p_tbl.columns[0].width = Inches(2.2)
        for c in range(1, cols):
            p_tbl.columns[c].width = Inches(3.8 / len(metrics))

        _cell(p_tbl, 0, 0, "Company", bold=True, color=_WHITE, fill=_NAVY, size=10)
        for c, lbl in enumerate(labels, 1):
            _cell(p_tbl, 0, c, lbl, bold=True, color=_WHITE, fill=_NAVY, size=10, align=PP_ALIGN.CENTER)

        def _peer_val(d, k):
            v = d.get(k)
            return f"{v:.1f}" if v is not None else "—"

        _cell(p_tbl, 1, 0, f"{tgt.get('name', stats.get('ticker',''))} ◀",
              bold=True, color=_WHITE, fill=_STEEL, size=10)
        for c, m in enumerate(metrics, 1):
            _cell(p_tbl, 1, c, _peer_val(tgt, m),
                  bold=True, color=_WHITE, fill=_STEEL, size=10, align=PP_ALIGN.CENTER)

        for r, peer in enumerate(peers, 2):
            rf = _LGREY if r % 2 == 0 else _WHITE
            _cell(p_tbl, r, 0, peer.get("name", peer.get("ticker", "—")), color=_DGREY, fill=rf, size=10)
            for c, m in enumerate(metrics, 1):
                _cell(p_tbl, r, c, _peer_val(peer, m), color=_DGREY, fill=rf, size=10, align=PP_ALIGN.CENTER)

    # Right: peer median callouts
    _txt(slide, "Peer Medians", Inches(6.8), _CT, Inches(6.1), Inches(0.3),
         size=13, bold=True, color=_NAVY)
    median_items = [
        ("Fwd P/E",     f"{_sf(medians.get('fpe'),1):.1f}x" if _sf(medians.get('fpe')) else "—"),
        ("Rev Growth",  f"{_sf(medians.get('rev_growth'),1):.1f}%" if _sf(medians.get('rev_growth')) else "—"),
        ("Op Margin",   f"{_sf(medians.get('op_margin'),1):.1f}%" if _sf(medians.get('op_margin')) else "—"),
        ("ROE",         f"{_sf(medians.get('roe'),1):.1f}%" if _sf(medians.get('roe')) else "—"),
    ]
    for i, (lbl, val) in enumerate(median_items):
        bx = Inches(6.8) + Inches(i * 1.55)
        _rect(slide, bx, Inches(1.1), Inches(1.4), Inches(0.85), fill=_LGREY, border=_LGREY)
        _txt(slide, val, bx, Inches(1.1), Inches(1.4), Inches(0.55),
             size=18, bold=True, color=_NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _txt(slide, lbl, bx, Inches(1.65), Inches(1.4), Inches(0.25),
             size=9, color=_MGREY, align=PP_ALIGN.CENTER)

    # Source
    source = comp_result.get("source", "") if ok else ""
    if source:
        _txt(slide, f"Source: yfinance {source} universe",
             _ML, Inches(5.5), Inches(12.5), Inches(0.25),
             size=9, italic=True, color=_MGREY)


def _slide_coverage(prs, cov_result: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Analyst Coverage", 10)

    ok      = cov_result and not cov_result.get("error")
    rating  = _sa(cov_result, "consensus_rating") if ok else "—"
    bull    = _sf(_sa(cov_result, "bull_ratio")) if ok else None
    buys    = _sa(cov_result, "buy_count",  default=0) if ok else 0
    holds   = _sa(cov_result, "hold_count", default=0) if ok else 0
    sells   = _sa(cov_result, "sell_count", default=0) if ok else 0
    total   = _sa(cov_result, "total_analysts", default=0) if ok else 0
    mean_t  = _sa(cov_result, "mean_target") if ok else None
    hi_t    = _sa(cov_result, "high_target") if ok else None
    lo_t    = _sa(cov_result, "low_target")  if ok else None
    upside  = _sa(cov_result, "upside_pct")  if ok else None
    recents = cov_result.get("recent_targets", [])[:4] if ok else []

    # Left panel: consensus
    rc = _rec_color(rating)
    _rect(slide, _ML, _CT, Inches(3.5), Inches(0.85), fill=rc)
    _txt(slide, rating, _ML, _CT, Inches(3.5), Inches(0.85),
         size=28, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _txt(slide, f"{total} analysts covering", _ML, Inches(1.67), Inches(3.5), Inches(0.3),
         size=11, color=_MGREY, align=PP_ALIGN.CENTER)

    # Buy/hold/sell bar
    if total and total > 0:
        bar_top    = Inches(2.1)
        bar_h      = Inches(0.32)
        bar_w      = Inches(3.5)
        dist_total = buys + holds + sells
        if dist_total > 0:
            seg_x = _ML
            for count, fill in [(buys, _GREEN), (holds, _GOLD), (sells, _RED)]:
                seg_w = bar_w * (count / dist_total)
                if seg_w > 0:
                    _rect(slide, seg_x, bar_top, seg_w, bar_h, fill=fill)
                    seg_x = seg_x + seg_w
            _txt(slide, f"Buy: {buys}", _ML, bar_top + Inches(0.35), Inches(1.2), Inches(0.25),
                 size=10, bold=True, color=_GREEN)
            _txt(slide, f"Hold: {holds}", _ML + Inches(1.2), bar_top + Inches(0.35), Inches(1.1), Inches(0.25),
                 size=10, bold=True, color=_GOLD)
            _txt(slide, f"Sell: {sells}", _ML + Inches(2.4), bar_top + Inches(0.35), Inches(1.1), Inches(0.25),
                 size=10, bold=True, color=_RED)
            if bull is not None:
                _txt(slide, f"Bull ratio: {bull:.1f}%",
                     _ML, Inches(2.7), Inches(3.5), Inches(0.3),
                     size=11, bold=True, color=_GREEN if bull >= 60 else _MGREY, align=PP_ALIGN.CENTER)
        else:
            _rect(slide, _ML, bar_top, bar_w, bar_h, fill=_rec_color(rating))
            _txt(slide, f"{total} analysts  |  Breakdown unavailable",
                 _ML, bar_top + Inches(0.35), bar_w, Inches(0.25),
                 size=10, italic=True, color=_MGREY)

    # Center panel: price targets
    _txt(slide, "Price Target Summary", Inches(4.2), _CT, Inches(4.5), Inches(0.3),
         size=13, bold=True, color=_NAVY)
    pt_items = [
        ("Mean Target",    _fmt_price(mean_t)),
        ("High Target",    _fmt_price(hi_t)),
        ("Low Target",     _fmt_price(lo_t)),
        ("Implied Upside", _fmt_pct(upside, plus=True)),
    ]
    y_pt = Inches(1.1)
    for label, value in pt_items:
        _txt(slide, label,  Inches(4.2), y_pt, Inches(2.3), Inches(0.35), size=10, bold=True, color=_NAVY)
        val_color = _GREEN if "+" in value else (_RED if "−" in value or (value != "—" and "-" in value) else _DGREY)
        _txt(slide, value, Inches(6.6), y_pt, Inches(2.0), Inches(0.35), size=10, color=val_color)
        y_pt += Inches(0.38)

    # Right panel: recent targets
    _txt(slide, "Recent Price Targets", Inches(9.0), _CT, Inches(4.0), Inches(0.3),
         size=13, bold=True, color=_NAVY)
    if recents:
        r_shp = slide.shapes.add_table(
            len(recents) + 1, 3, Inches(9.0), Inches(1.1), Inches(4.0), Inches(1.8))
        r_tbl = r_shp.table
        r_tbl.columns[0].width = Inches(1.9)
        r_tbl.columns[1].width = Inches(1.1)
        r_tbl.columns[2].width = Inches(1.0)
        for c, h in enumerate(["Firm", "Target", "Date"]):
            _cell(r_tbl, 0, c, h, bold=True, color=_WHITE, fill=_NAVY, size=10, align=PP_ALIGN.CENTER)
        for r, rec in enumerate(recents, 1):
            rf = _LGREY if r % 2 == 0 else _WHITE
            _cell(r_tbl, r, 0, rec.get("firm", "—")[:20], color=_DGREY, fill=rf, size=9)
            _cell(r_tbl, r, 1, _fmt_price(rec.get("price_target")), color=_DGREY, fill=rf, size=9, align=PP_ALIGN.CENTER)
            _cell(r_tbl, r, 2, rec.get("date", "—")[:10], color=_DGREY, fill=rf, size=9, align=PP_ALIGN.CENTER)
    else:
        _txt(slide, "No recent target data.", Inches(9.0), Inches(1.1), Inches(4.0), Inches(0.4),
             size=10, color=_MGREY, italic=True)


def _slide_risks(prs, research: dict | None,
                 stats: dict | None = None, fin_data: dict | None = None,
                 dcf_result: dict | None = None, comp_result: dict | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Key Risks", 11)

    thesis = (research or {}).get("thesis") or {}
    ok     = not thesis.get("_placeholder")
    bears  = thesis.get("bear", []) if ok else []

    if len(bears) >= 2:
        risks = bears[:3]
    elif stats and fin_data:
        risks = _build_data_risks(stats, fin_data, dcf_result, comp_result)
    else:
        risks = [
            "Macro / rate sensitivity — multiple compression risk if rates stay elevated",
            "Competitive disruption — faster-than-expected share loss to new entrants",
            "Execution risk — failure to meet consensus revenue or margin guidance",
        ]
    probs = ["High", "Medium", "Medium"]

    positions = [_ML, Inches(4.65), Inches(8.9)]
    card_w    = Inches(4.0)

    for i, (text, prob) in enumerate(zip(risks, probs)):
        cx = positions[i]
        _rect(slide, cx, _CT, card_w, Inches(4.2), fill=_LGREY, border=_LGREY)
        _rect(slide, cx, _CT, card_w, Inches(0.06), fill=_RED)

        prob_color = _RED if prob == "High" else (_GOLD if prob == "Medium" else _GREEN)
        _rect(slide, cx, _CT + Inches(0.1), Inches(1.1), Inches(0.3), fill=prob_color)
        _txt(slide, prob, cx, _CT + Inches(0.1), Inches(1.1), Inches(0.3),
             size=10, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        display = text if len(text) <= 250 else text[:247] + "…"
        _txt(slide, display,
             cx + Inches(0.1), _CT + Inches(0.5),
             card_w - Inches(0.2), Inches(3.5),
             size=11, color=_DGREY, wrap=True)



def _slide_final(prs, stats: dict, research: dict | None, cov_result: dict | None,
                 transcript_result: dict | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, _W, _H, fill=_NAVY)

    _txt(slide, _GS_LBL, Inches(9.5), Inches(0.15), Inches(3.5), Inches(0.3),
         size=10, italic=True, color=_MGREY, align=PP_ALIGN.RIGHT)

    thesis  = (research or {}).get("thesis") or {}
    ok      = not thesis.get("_placeholder")
    rating  = _fallback_rating(research, cov_result)
    target  = _fallback_target(research, cov_result)
    verdict = thesis.get("verdict", "") if ok else ""
    bulls   = thesis.get("bull", []) if ok else []

    ticker = stats.get("ticker", "")
    px     = _fmt_price(stats.get("current_price"))
    upside = _fmt_pct(_sa(cov_result, "upside_pct"), plus=True) if cov_result and not cov_result.get("error") else "—"

    rc = _rec_color(rating)
    _txt(slide, rating,
         Inches(1.0), Inches(0.9), Inches(11.3), Inches(1.3),
         size=72, bold=True, color=rc, align=PP_ALIGN.CENTER)

    _txt(slide, f"{ticker}  |  Target: {target}  |  Upside: {upside}  |  12-Month Horizon",
         Inches(1.0), Inches(2.3), Inches(11.3), Inches(0.45),
         size=18, color=_WHITE, align=PP_ALIGN.CENTER)

    if verdict:
        _txt(slide, verdict, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.7),
             size=14, italic=True, color=_LGREY, align=PP_ALIGN.CENTER, wrap=True)

    # Conviction — use shared utility
    bull_ratio_val = cov_result.get("bull_ratio") if (cov_result and not cov_result.get("error")) else None
    beat_streak_val = (transcript_result.get("beat_streak", 0)
                       if (transcript_result and not transcript_result.get("error")) else 0)
    conviction = get_conviction(bull_ratio_val, beat_streak_val)
    conv_color = _GREEN if conviction == "High" else (_RED if conviction == "Low" else _GOLD)
    _txt(slide, f"Conviction: {conviction.upper()}",
         Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.4),
         size=16, bold=True, color=conv_color, align=PP_ALIGN.CENTER)

    if bulls:
        summary_pts = "  ·  ".join(b[:55] for b in bulls[:3])
        _txt(slide, summary_pts,
             Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.55),
             size=10, color=_MGREY, align=PP_ALIGN.CENTER, wrap=True)

    _txt(slide, _GS_LBL,
         Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.25),
         size=10, italic=True, color=_MGREY, align=PP_ALIGN.CENTER)

    _txt(slide, "12", Inches(12.7), Inches(7.2), Inches(0.35), Inches(0.2),
         size=9, color=_MGREY, align=PP_ALIGN.RIGHT)


# ── Entry point ───────────────────────────────────────────────────────────────

def run_pitch(ticker: str, stats: dict, fin_data: dict,
              dcf_result: dict | None = None,
              research: dict | None = None,
              comp_result: dict | None = None,
              cov_result: dict | None = None,
              transcript_result: dict | None = None,
              out_path: str = "") -> dict:
    """Build a 12-slide UC Lindner pitch deck. Never raises."""
    try:
        prs = Presentation()
        prs.slide_width  = _W
        prs.slide_height = _H

        _slide_cover(prs, stats, research, cov_result)
        _slide_summary(prs, stats, fin_data, research, comp_result, cov_result, dcf_result)
        _slide_overview(prs, stats, fin_data)
        _slide_thesis(prs, research, stats, fin_data, comp_result)
        _slide_financials(prs, stats, fin_data)
        _slide_dcf(prs, stats, fin_data, dcf_result)
        _slide_comps(prs, stats, research, comp_result)
        _slide_football(prs, stats, dcf_result, comp_result, cov_result)
        _slide_competitive(prs, stats, comp_result)
        _slide_coverage(prs, cov_result)
        _slide_risks(prs, research, stats, fin_data, dcf_result, comp_result)
        _slide_final(prs, stats, research, cov_result, transcript_result=transcript_result)

        prs.save(out_path)
        return {"error": None, "path": out_path, "n_slides": 12}

    except Exception as exc:
        return {"error": str(exc)}
