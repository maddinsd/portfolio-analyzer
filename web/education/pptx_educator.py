"""
PowerPoint educator — adds speaker notes to an existing presentation.
Matches slides by index (1-12). Never regenerates slide content.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Pt


def add_ppt_notes(pptx_path: str, ppt_notes: list) -> int:
    """
    Adds speaker notes to each slide in pptx_path.
    ppt_notes: list of {"slide": int, "title": str, "notes": str}
    Returns number of slides updated.
    """
    if not ppt_notes:
        return 0

    prs = Presentation(pptx_path)
    notes_by_index = {n["slide"]: n.get("notes", "") for n in ppt_notes if n.get("slide")}
    slides_updated = 0

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        note_text = notes_by_index.get(slide_num, "")
        if not note_text:
            continue

        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame

        # Clear existing notes and set new text
        tf.clear()
        tf.text = note_text

        # Set a readable font size for notes
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)

        slides_updated += 1

    prs.save(pptx_path)
    return slides_updated
